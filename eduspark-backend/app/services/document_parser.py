"""文档解析服务 — 提取 PDF/PPT/Word 文本内容"""
import re
from pathlib import Path


class DocumentParser:
    """支持 PDF、PPT、Word 的文本提取"""

    @staticmethod
    def parse(file_path: str) -> dict:
        """
        自动识别文件类型并提取文本
        返回: {"success": bool, "text": str, "pages": int, "file_type": str, "error": str}
        """
        path = Path(file_path)
        suffix = path.suffix.lower()

        if not path.exists():
            return {"success": False, "error": "文件不存在"}

        try:
            if suffix == ".pdf":
                return DocumentParser._parse_pdf(file_path)
            elif suffix in (".ppt", ".pptx"):
                return DocumentParser._parse_ppt(file_path)
            elif suffix in (".doc", ".docx"):
                return DocumentParser._parse_docx(file_path)
            elif suffix in (".txt", ".md"):
                return DocumentParser._parse_text(file_path)
            else:
                return {"success": False, "error": f"不支持的文件格式: {suffix}"}
        except Exception as e:
            return {"success": False, "error": str(e), "file_type": suffix}

    @staticmethod
    def _parse_pdf(file_path: str) -> dict:
        """解析 PDF 文件"""
        text_parts = []
        page_count = 0

        # 方法1: pymupdf（首选，中文支持好）
        try:
            import fitz
            doc = fitz.open(file_path)
            page_count = len(doc)
            for page in doc:
                text = page.get_text()
                if text.strip():
                    text_parts.append(text.strip())
            doc.close()
            if text_parts:
                return {
                    "success": True,
                    "text": "\n\n".join(text_parts),
                    "pages": page_count,
                    "file_type": "pdf",
                }
        except ImportError:
            pass

        # 方法2: pdfplumber（备选）
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        text_parts.append(text.strip())
            if text_parts:
                return {
                    "success": True,
                    "text": "\n\n".join(text_parts),
                    "pages": page_count,
                    "file_type": "pdf",
                }
        except ImportError:
            pass

        return {"success": False, "error": "PDF解析需要安装 pymupdf 或 pdfplumber"}

    @staticmethod
    def _parse_ppt(file_path: str) -> dict:
        """解析 PPT 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_parts.append(text)
                if slide_parts:
                    slides_text.append(f"幻灯片 {i + 1}:\n" + "\n".join(slide_parts))
            return {
                "success": True,
                "text": "\n\n".join(slides_text),
                "pages": len(slides_text),
                "file_type": "ppt",
            }
        except ImportError:
            return {"success": False, "error": "PPT解析需要安装 python-pptx"}

    @staticmethod
    def _parse_docx(file_path: str) -> dict:
        """解析 Word 文件"""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            return {
                "success": True,
                "text": "\n\n".join(paragraphs),
                "pages": len(paragraphs),
                "file_type": "docx",
            }
        except ImportError:
            return {"success": False, "error": "Word解析需要安装 python-docx"}

    @staticmethod
    def _parse_text(file_path: str) -> dict:
        """解析纯文本文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        lines = text.strip().split("\n")
        return {
            "success": True,
            "text": text.strip(),
            "pages": len(lines),
            "file_type": "txt",
        }

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[dict]:
        """
        将文本按段落分块，用于后续向量化
        返回: [{"index": 0, "content": "...", "char_count": 800}, ...]
        """
        # 按段落分割
        paragraphs = re.split(r'\n\n+', text)
        chunks = []
        current_chunk = ""
        chunk_index = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # 如果当前块+新段落超过限制，保存当前块
            if len(current_chunk) + len(para) > chunk_size and current_chunk:
                chunks.append({
                    "index": chunk_index,
                    "content": current_chunk.strip(),
                    "char_count": len(current_chunk.strip()),
                })
                chunk_index += 1
                # 保留 overlap 部分
                if len(current_chunk) > overlap:
                    current_chunk = current_chunk[-overlap:] + "\n\n" + para
                else:
                    current_chunk = para
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para

        # 最后一个块
        if current_chunk.strip():
            chunks.append({
                "index": chunk_index,
                "content": current_chunk.strip(),
                "char_count": len(current_chunk.strip()),
            })

        return chunks


# 全局单例
doc_parser = DocumentParser()
